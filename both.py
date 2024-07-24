import streamlit as st
import pandas as pd
import openai
from pptx import Presentation
import io
from builtins import str
import zipfile
import os

# Initialize OpenAI API
openai.api_key = ""


def generate_topics(industry):
    prompt = (
        f"Generate 3 distinct book topics and 3 distinct LinkedIn post topic ideas for the {industry} industry. "
        f"Ensure the topics are general and not related to gender, women in leadership, or personal branding. Format the response as:\n"
        f"1. Book Topic 1\n2. Book Topic 2\n3. Book Topic 3\n4. LinkedIn Post Idea 1\n5. LinkedIn Post Idea 2\n6. LinkedIn Post Idea 3"
    )
    response = openai.Completion.create(
        engine="gpt-3.5-turbo-instruct", prompt=prompt, max_tokens=150
    )
    lines = [
        line.strip()
        for line in response.choices[0].text.strip().split("\n")
        if line.strip()
    ]
    if len(lines) == 6:
        return lines
    else:
        return []


def apply_text_replacements(prs, replacements):
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)


def replace_text_in_shape(shape, replacements):
    if shape.has_text_frame:
        text_frame = shape.text_frame
        while True:
            updated = False
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    for key, value in replacements.items():
                        if key in original_text:
                            run.text = original_text.replace(key, value)
                            updated = True
                    if original_text != run.text:
                        updated = True
            if not updated:
                break
    if shape.shape_type == 6:  # Group shape
        for sub_shape in shape.shapes:
            replace_text_in_shape(sub_shape, replacements)


# Streamlit UI
st.title("Industry-specific Topics Generator and PowerPoint Integrator")

# Upload CSV file
csv_file = st.file_uploader("Upload a CSV file", type="csv")

# Upload PowerPoint template
pptx_file = st.file_uploader("Upload your PPTX template", type="pptx")

if csv_file and pptx_file:
    df = pd.read_csv(csv_file)
    template_bytes = pptx_file.read()

    if "Industry" in df.columns:
        for index, row in df.iterrows():
            industry = row["Industry"]
            topics = generate_topics(industry)

            retries = 0
            while len(topics) < 6 and retries < 3:
                topics = generate_topics(industry)
                retries += 1

            if len(topics) == 6:
                for i in range(3):
                    df.at[index, f"Topic{i+1}"] = (
                        topics[i].split(". ", 1)[-1].replace('"', "")
                    )
                    df.at[index, f"Post{i+1}"] = (
                        topics[i + 3].split(". ", 1)[-1].replace('"', "")
                    )

        with zipfile.ZipFile("presentations.zip", "w") as zipf:
            for index, row in df.iterrows():
                prs = Presentation(io.BytesIO(template_bytes))
                replacements = {f"{{{{{col}}}}}": str(row[col]) for col in df.columns}
                apply_text_replacements(prs, replacements)

                output_file = f"{row['First Name']} - Thought Leader.pptx"
                prs.save(output_file)
                zipf.write(output_file, arcname=output_file)
                os.remove(output_file)

        st.success("PowerPoints updated and zipped successfully!")
        with open("presentations.zip", "rb") as zip_file:
            st.download_button(
                label="Download All Presentations",
                data=zip_file,
                file_name="presentations.zip",
                mime="application/zip",
            )
    else:
        st.error(
            "The CSV file does not contain the 'Industry' column. Please check the file and try again."
        )
