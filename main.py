from fastapi import FastAPI, Request, UploadFile, File
from fastapi.responses import JSONResponse
from loguru import logger
from db_models import *
from utils import (
    verify_password,
    generate_salt_and_hash,
    decrypt_password,
    encrypt_password,
)
from bucket_utils import upload_to_bucket,download_from_bucket
from email_handler import *
import json
from constants import GLOBAL_EMAIL_DECRYPTOR_KEY
from fastapi.background import BackgroundTasks
import random
from datetime import timedelta
import shutil
from mongoengine.errors import NotUniqueError
from datetime import datetime
from bson import ObjectId
import uuid
from builtins import str
from fastapi import FastAPI, BackgroundTasks, HTTPException
from pydantic import BaseModel, EmailStr
from apscheduler.schedulers.background import BackgroundScheduler
from datetime import datetime
import smtplib
from email.mime.text import MIMEText
import pandas as pd
import openai
from pptx import Presentation
import io
import os


scheduler = BackgroundScheduler()
scheduler.start()

app = FastAPI()

sign_up_otp_handler = SignupOTP()
user_data_manage = UserDataManage()


@app.get("/")
def healthcheck():
    return {"status": "healthy"}


@app.post("/login")
async def login(request: Request):
    logger.info("login api called")
    data = await request.json()
    doc = {"email_id": data.get("email_id"), "password": data.get("password")}
    try:
        logger.info("login api started to execute")
        user = users_collection.find_one({"email_id": doc["email_id"]})

        if user is None:
            return {"status": False, "message": "please check your email/password !"}
        else:
            users_collection.update_one(
                {"email_id": doc["email_id"]}, {"$set": {"user_id": str(user["_id"])}}
            )
            users_data = users_collection.find_one({"email_id": doc["email_id"]})
            password = doc["password"]
            is_pwd_correct = verify_password(
                password=password,
                stored_hash=users_data["password"],
                salt=users_data["salted_pwd"],
            )
            if is_pwd_correct:
                user_id = users_data["user_id"]
                login_session_data = login_session_collection.find_one(
                    {"email_id": doc["email_id"], "user_id": user_id}
                )
                if login_session_data is None:
                    login_session = LoginSession(
                        email_id=doc["email_id"], user_id=user_id
                    )
                    login_session_collection.insert_one(login_session.__dict__)
                else:
                    login_session_collection.update_one(
                        {"email_id": doc["email_id"], "user_id": user_id, "status": 0},
                        {"$set": {"status": 1}},
                    )
                if not user_data_manage.is_user_id_in_global_metadata(user_id=user_id):
                    user_data_manage.insert_session_id(user_id=user_id)
                response = {
                    "status": True,
                    "message": "Login successful",
                    "user_id": user_id,
                    "email_id": doc["email_id"],
                    "username": users_data["first_name"]
                    + " "
                    + users_data["last_name"],
                }
                logger.info(response)
                return response
            else:
                return {"message": "Please check your password..."}
    except Exception as e:
        logger.info(f"Error: {str(e)}")
        return {
            "status": False,
            "message": "something went wrong!",
            "message": str(e),
        }


@app.post("/sign_up")
async def sign_up(request: Request, backgroud_task: BackgroundTasks):
    data = await request.json()
    doc = {
        "first_name": data.get("first_name"),
        "last_name": data.get("last_name"),
        "email_id": data.get("email_id"),
        "password": data.get("password"),
    }
    try:
        is_valid = check_email(email=doc["email_id"])
        if is_valid:
            user_exists = users_collection.find_one({"email_id": doc["email_id"]})
            if user_exists is not None:
                return {"status": False, "message": "User already exists"}
            otp = str(random.randint(100000, 999999))
            if sign_up_otp_handler.is_in_otp_data(doc["email_id"]):
                sign_up_otp_handler.update_otp_data(doc["email_id"], otp)
            else:
                sign_up_otp_handler.insert_otp_data(doc["email_id"], otp)
            logger.info("Sending otp to user")
            email_subject, email_body = email_template_otp(
                user_name=doc["first_name"], otp_code=otp
            )
            send_otp(email_body, email_subject, doc["email_id"], backgroud_task)
            logger.success("OTP sent successfully")
            return {"status": True, "message": "Otp sent successful please verify"}
        return {"status": False, "message": "Invalid Email Address Provided"}
    except NotUniqueError as e:
        print(str(e))
        logger.info(f"Went wrong in signup {str(e)}")
        return {
            "status": False,
            "message": "user found with same email id",
            "error": str(e),
        }
    except Exception as e:
        print(str(e))
        return {
            "status": False,
            "message": "unknown error",
            "error": str(e),
        }


@app.post("/verify_otp")
async def verify_otp(request: Request):
    data = await request.json()
    otp_user_input = {
        "first_name": data.get("first_name"),
        "last_name": data.get("last_name"),
        "email_id": data.get("email_id"),
        "password": data.get("password"),
        "otp": data.get("otp"),
    }
    try:
        if sign_up_otp_handler.is_in_otp_data(otp_user_input["email_id"]):
            otp, timestamp = sign_up_otp_handler.get_otp_data(
                otp_user_input["email_id"]
            )
            current_time = datetime.now()
            if current_time - timestamp <= timedelta(minutes=5):
                if otp_user_input["otp"] == otp:
                    sign_up_otp_handler.delete_otp_data(otp_user_input["email_id"])
                    salt, password = generate_salt_and_hash(otp_user_input["password"])
                    user = User(
                        first_name=otp_user_input["first_name"],
                        last_name=otp_user_input["last_name"],
                        email_id=otp_user_input["email_id"],
                        password=password,
                        salted_pwd=salt,
                    )
                    users_collection.insert_one(user.__dict__)
                    logger.success(
                        'successfully verified otp with status "status": True, "message": "successfully verified otp"'
                    )
                    return {"status": True, "message": "successfully verified otp"}
                else:
                    logger.warning("OTP verification failed, please check and re-enter")
                    return {
                        "status": False,
                        "message": "OTP verification failed, please check and re-enter",
                    }
            else:
                return {
                    "status": False,
                    "message": "OTP has expired please generate new one...",
                }
        else:
            logger.warning("OTP either expired or invalid, please generate a new one")
            return {
                "status": False,
                "message": "OTP either expired or invalid, please generate a new one",
            }
    except Exception as e:
        logger.error(f'"status": "False reason" {str(e)}"')
        return {"status": False, "reason": str(e)}


class EmailRequest(BaseModel):
    sender_email: EmailStr
    subject: str
    body: str
    user_id: str
    doc_id: str
    send_at: datetime


class ConnectSenderRequest(BaseModel):
    user_id: str
    sender_email: EmailStr
    password: str


@app.post("/connect_sender")
async def connect_sender(request: ConnectSenderRequest):
    already_exists = users_collection_email_scheduler.find_one(
        {"sender_email": request.sender_email, "user_id": request.user_id}
    )
    if already_exists:
        return {"status": "false", "message": "Sender email already connected"}
    encrypted_password = encrypt_password(
        password=request.password, key=GLOBAL_EMAIL_DECRYPTOR_KEY
    )
    user_data = {
        "user_id": request.user_id,
        "sender_email": request.sender_email,
        "password": encrypted_password,
    }
    users_collection_email_scheduler.update_one(
        {"user_id": request.user_id, "sender_email": request.sender_email},
        {"$set": user_data},
        upsert=True,
    )
    return {"message": "Sender email connected successfully"}


@app.post("/schedule_email")
async def schedule_email(email_request: EmailRequest):
    user = users_collection_email_scheduler.find_one(
        {"user_id": email_request.user_id, "sender_email": email_request.sender_email}
    )
    if not user:
        raise HTTPException(status_code=404, detail="Sender email not connected")
    get_docs = user_data_manage.get_document(
        user_id=email_request.user_id, doc_id=email_request.doc_id
    )
    if "email_id" not in get_docs["dataframe_fields"]:
        return {
            "status": "false",
            "message": "Cannot send email please check the field email_id exists in the csv.",
        }
    decrypted_pasword = decrypt_password(
        encrypted_password=user["password"], key=GLOBAL_EMAIL_DECRYPTOR_KEY
    )
    job_data = email_request.dict()
    job_data["user_id"] = email_request.user_id
    job_data["status"] = "scheduled"
    job_id = jobs_collection_email_schedule.insert_one(job_data).inserted_id
    scheduler.add_job(
        send_email,
        "date",
        run_date=email_request.send_at,
        args=[
            str(job_id),
            decrypted_pasword,
            email_request.doc_id,
            get_docs["dataframe"],
        ],
        id=str(job_id),
    )
    return {"message": "Email scheduled successfully"}

@app.get("/get_connected_emails")
async def get_connected_emails(user_id: str):
    projection={"sender_email":1,"_id":0}
    connected_emails = list(users_collection_email_scheduler.find({"user_id":user_id},projection=projection))
    return connected_emails


@app.get("/get_email_jobs")
async def get_email_jobs(user_id: str):
    jobs = list(jobs_collection_email_schedule.find({"user_id":user_id}))
    for job in jobs:
        job["_id"] = str(job["_id"])
    return jobs


def send_email(job_id: str, password: str, doc_id: str, get_docs):
    job = jobs_collection_email_schedule.find_one({"_id": ObjectId(job_id)})
    if not job:
        print(f"Job {job_id} not found")
        return {"status": "false", "message": "no such email job found"}
    if not get_docs:
        print(f"document {doc_id} not found")
        return {"status": "false", "message": "no such document found"}
    for data in get_docs:
        if data["email_id"] == "":
            continue
        sender_email: str = job["sender_email"]
        receiver_email: str = data["email_id"]
        subject: str = job["subject"].format(**data)
        body: str = job["body"].format(**data)

        msg = MIMEText(body)
        msg["Subject"] = subject
        msg["From"] = sender_email
        msg["To"] = receiver_email
        try:
            with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
                server.login(sender_email, password)
                server.sendmail(sender_email, receiver_email, msg.as_string())
            jobs_collection_email_schedule.update_one(
                {"_id": ObjectId(job_id)}, {"$set": {"status": "sent"}}
            )
            print(f"Email sent from {doc_id} sender: {sender_email} to reciever:{receiver_email}")
        except Exception as e:
            jobs_collection_email_schedule.update_one(
                {"_id": ObjectId(job_id)},
                {"$set": {"status": "failed", "reason": str(e)}},
            )
            print(f"Failed to send email: {e}")


def generate_topics(
    industry,
    prompt_template,
    openai_api_key="",
    gpt_model="gpt-3.5-turbo-instruct",
):
    openai.api_key = openai_api_key
    if prompt_template is None:
        prompt = (
            f"Generate 3 distinct book topics and 3 distinct LinkedIn post topic ideas for the {industry} industry. "
            f"Ensure the topics are general and not related to gender, women in leadership, or personal branding. Format the response as:\n"
            f"1. Book Topic 1\n2. Book Topic 2\n3. Book Topic 3\n4. LinkedIn Post Idea 1\n5. LinkedIn Post Idea 2\n6. LinkedIn Post Idea 3"
        )
    response = openai.Completion.create(engine=gpt_model, prompt=prompt, max_tokens=150)
    lines = [
        line.strip()
        for line in response.choices[0].text.strip().split("\n")
        if line.strip()
    ]
    if len(lines) == 6:
        return lines
    else:
        return []


def apply_text_replacements(prs, replacements):
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)


def replace_text_in_shape(shape, replacements):
    if shape.has_text_frame:
        text_frame = shape.text_frame
        while True:
            updated = False
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    for key, value in replacements.items():
                        if key in original_text:
                            run.text = original_text.replace(key, value)
                            updated = True
                    if original_text != run.text:
                        updated = True
            if not updated:
                break
    if shape.shape_type == 6:
        for sub_shape in shape.shapes:
            replace_text_in_shape(sub_shape, replacements)


@app.post("/process_data")
async def upload_files(
    user_id: str,
    doc_id: str,
    openai_api_key: str,
    gpt_model: str,
    topic: str,
    prompt_template=None
):
    try:
        if user_data_manage.is_user_id_in_global_metadata(user_id=user_id):
            get_document = user_data_manage.get_document(user_id=user_id, doc_id=doc_id)
            template_bytes = download_from_bucket(user_id=user_id,doc_id=doc_id,type="user",extension="pptx")
            df = pd.DataFrame(get_document["dataframe"])
            if topic not in df.columns:
                raise HTTPException(
                    status_code=400,
                    detail=f"The CSV file does not contain the {topic} column.",
                )
            for index, row in df.iterrows():
                industry = row[topic]
                topics = generate_topics(
                    industry,
                    prompt_template=prompt_template,
                    openai_api_key=openai_api_key,
                    gpt_model=gpt_model,
                )
                retries = 0
                while len(topics) < 6 and retries < 3:
                    topics = generate_topics(
                        industry,
                        prompt_template=prompt_template,
                        openai_api_key=openai_api_key,
                        gpt_model=gpt_model,
                    )
                    retries += 1
                if len(topics) == 6:
                    for i in range(3):
                        df.at[index, f"Topic{i+1}"] = (
                            topics[i].split(". ", 1)[-1].replace('"', "")
                        )
                        df.at[index, f"Post{i+1}"] = (
                            topics[i + 3].split(". ", 1)[-1].replace('"', "")
                        )
            os.makedirs(user_id, exist_ok=True)
            for index, row in df.iterrows():
                prs = Presentation(io.BytesIO(template_bytes))
                replacements = {f"{{{{{col}}}}}": str(row[col]) for col in df.columns}
                apply_text_replacements(prs, replacements)
                output_file = f"{doc_id}_presentation.pptx"
                prs_path = os.path.join(user_id, output_file)
                prs.save(prs_path)
            df.to_csv(f"{user_id}/{doc_id}.csv")
            upload_to_bucket(
                f"{user_id}/{doc_id}_presentation.pptx",user_id=user_id, doc_id=doc_id, type="modified",extension="pptx"
            )
            upload_to_bucket(
                f"{user_id}/{doc_id}.csv", user_id=user_id, doc_id=doc_id, type="modified",extension="csv"
            )
            download_ppt_link = f"https://storage.googleapis.com/bucketsuhas/{user_id}/{doc_id}_modified_pptx.pptx"
            download_csv_link = f"https://storage.googleapis.com/bucketsuhas/{user_id}/{doc_id}_modified_csv.csv"
            user_data_manage.update_document_collection(
                user_id=user_id,
                document_id=doc_id,
                prompt_template=prompt_template,
                topics=topic,
                download_ppt=download_ppt_link,
                download_csv_link=download_csv_link,
            )
            return JSONResponse(
                {
                    "status": "true",
                    "download_ppt_link": download_ppt_link,
                    "download_csv_link": download_csv_link,
                    "topic_selected": topic,
                }
            )
        return JSONResponse({"status": "false", "message": "No such user found"})
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})
    finally:
        if os.path.exists(user_id):
            shutil.rmtree(user_id)


class AvaialableField(BaseModel):
    user_id: str


@app.get("/get_document")
def show_document(user_id: str):
    documents = user_data_manage.get_all_document(user_id=user_id)
    print(documents)
    return documents[user_id]


@app.post("/upload_ppt_csv")
async def available_fields(
    user_id: str, csv_file: UploadFile = File(...), pptx_file: UploadFile = File(...)
):
    if user_data_manage.is_user_id_in_global_metadata(user_id=user_id):
        df = pd.read_csv(csv_file.file)
        df.fillna("", inplace=True)
        df_dict = json.loads(json.dumps(df.to_dict(orient="records"), allow_nan=True))
        if len(df_dict) < 5:
            return {
                "status": "false",
                "message": "dataframe uploaded must have atleast 5 rows",
            }
        template_bytes = pptx_file.file.read()
        doc_id = str(uuid.uuid1())
        with open(f"{doc_id}.pptx", "wb") as f:
            f.write(template_bytes)
        upload_to_bucket(
            local=f"{doc_id}.pptx", doc_id=doc_id, user_id=user_id, type="user",extension="pptx"
        )
        sharable_link_ppt=f"https://storage.googleapis.com/bucketsuhas/{user_id}/{doc_id}_user_pptx.pptx"
        uploaded_time = user_data_manage.insert_documents(
            user_id=user_id,
            document_id=doc_id,
            template_name=str(pptx_file.filename),
            df_name=str(csv_file.filename),
            ppt_link=sharable_link_ppt,
            dataframe=df_dict,
            prompt_template=None,
            topics=None,
            dataframe_fields=df.columns.to_list(),
        )
        os.remove(f"{doc_id}.pptx")
        return {
            "status": "true",
            "doc_id": doc_id,
            "template_name": str(pptx_file.filename),
            "csv_name": str(csv_file.filename),
            "uploaded_time": uploaded_time,
            "ppt_link": sharable_link_ppt,
            "available_fields": df.columns.to_list(),
            "dataframe": df_dict,
        }
    return {"status": "false", "message": "No such user found"}


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
